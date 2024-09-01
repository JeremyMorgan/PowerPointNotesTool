from pptx import Presentation
import re

# Copyright (c) 2024 Jeremy Morgan
# This file is part of PowerPointNotesTool
#
# PowerPointNotesTool is free software: you can redistribute it and/or modify
# it under the terms of the GNU General Public License as published by
# the Free Software Foundation, either version 3 of the License, or
# (at your option) any later version.
#
# PowerPointNotesTool is distributed in the hope that it will be useful,
# but WITHOUT ANY WARRANTY; without even the implied warranty of
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
# GNU General Public License for more details.
#
# You should have received a copy of the GNU General Public License
# along with PowerPointNotesTool.  If not, see <https://www.gnu.org/licenses/>.

def extract_slide_notes(file_path):
    """
    Extracts slide notes from a text file formatted with specific patterns.

    The text file should contain slide notes enclosed within 'start slide <slide_number>' 
    and 'end slide <slide_number>' tags. Each slide's notes should be separated by these tags.

    Parameters:
    file_path (str): The path to the text file containing slide notes.

    Returns:
    dict: A dictionary where keys are slide numbers (int) and values are slide notes (str).
    """
    with open(file_path, 'r') as f:
        content = f.read()

    slide_pattern = r'start slide (\d+)(.*?)end slide \1'
    slides_notes = re.findall(slide_pattern, content, re.DOTALL)

    notes_dict = {}
    for slide_num, notes in slides_notes:
        notes_dict[int(slide_num)] = notes.strip()

    return notes_dict

# Open the PowerPoint file
prs = Presentation('presentation.pptx')

# Extract notes from text file
notes_dict = extract_slide_notes('notes.txt')

# Add notes to slides
for i, slide in enumerate(prs.slides, start=1):
    if i in notes_dict:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_dict[i]

# Save the modified presentation
prs.save('presentation_with_notes.pptx')

# Print a happy message
print("Notes added successfully!")